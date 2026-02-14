#!/usr/bin/env python3
import sys
import os
import time

def convert_pdf_to_powerpoint(input_path, output_path):
    try:
        print(f"Starting conversion: {input_path} -> {output_path}")
        start_time = time.time()

        import pymupdf as fitz
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        pdf_document = fitz.open(input_path)
        prs = Presentation()

        total_pages = len(pdf_document)

        for page_num in range(total_pages):
            print(f"Processing page {page_num + 1}/{total_pages}")
            page = pdf_document.load_page(page_num)

            page_width = page.rect.width
            page_height = page.rect.height

            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            slide_w = prs.slide_width
            slide_h = prs.slide_height
            scale_x = slide_w / page_width
            scale_y = slide_h / page_height

            tables = page.find_tables()
            table_areas = set()

            if tables and len(tables.tables) > 0:
                for table in tables.tables:
                    data = table.extract()
                    if not data or len(data) == 0:
                        continue

                    bbox = table.bbox
                    table_areas.add((round(bbox[1]), round(bbox[3])))

                    left = Emu(int(bbox[0] * scale_x))
                    top = Emu(int(bbox[1] * scale_y))
                    width = Emu(int((bbox[2] - bbox[0]) * scale_x))
                    height = Emu(int((bbox[3] - bbox[1]) * scale_y))

                    num_rows = len(data)
                    num_cols = max(len(row) for row in data) if data else 1

                    try:
                        tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
                        tbl = tbl_shape.table

                        for row_idx, row_data in enumerate(data):
                            for col_idx in range(num_cols):
                                cell = tbl.cell(row_idx, col_idx)
                                cell_text = row_data[col_idx] if col_idx < len(row_data) and row_data[col_idx] else ""
                                cell_text = cell_text.strip() if isinstance(cell_text, str) else str(cell_text) if cell_text else ""
                                cell.text = cell_text

                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(9)
                                        if row_idx == 0:
                                            run.font.bold = True
                    except Exception:
                        left = Emu(int(bbox[0] * scale_x))
                        top = Emu(int(bbox[1] * scale_y))
                        width = Emu(int((bbox[2] - bbox[0]) * scale_x))
                        height = Emu(int((bbox[3] - bbox[1]) * scale_y))
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        for row_data in data:
                            row_text = " | ".join([c.strip() if c else "" for c in row_data])
                            p = tf.add_paragraph()
                            p.text = row_text
                            p.font.size = Pt(9)

            blocks = page.get_text("dict")

            for block in blocks["blocks"]:
                if "lines" not in block:
                    continue

                block_y_top = block["bbox"][1]
                block_y_bottom = block["bbox"][3]

                in_table = False
                for (t_top, t_bottom) in table_areas:
                    if block_y_top >= t_top - 5 and block_y_bottom <= t_bottom + 5:
                        in_table = True
                        break
                if in_table:
                    continue

                bbox = block["bbox"]
                left = Emu(int(bbox[0] * scale_x))
                top = Emu(int(bbox[1] * scale_y))
                width = Emu(int((bbox[2] - bbox[0]) * scale_x))
                height = Emu(int((bbox[3] - bbox[1]) * scale_y))

                if width <= 0 or height <= 0:
                    continue

                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                first_para = True
                for line in block["lines"]:
                    line_text = ""
                    line_size = 11
                    line_bold = False
                    line_italic = False
                    line_color = None

                    for span in line["spans"]:
                        text = span["text"]
                        if text.strip():
                            line_text += text
                            line_size = span["size"]
                            if span["flags"] & 2**4:
                                line_bold = True
                            if span["flags"] & 2**1:
                                line_italic = True
                            color_int = span.get("color", 0)
                            if color_int and color_int != 0:
                                r = (color_int >> 16) & 0xFF
                                g = (color_int >> 8) & 0xFF
                                b = color_int & 0xFF
                                line_color = RGBColor(r, g, b)

                    if not line_text.strip():
                        continue

                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()

                    run = p.add_run()
                    run.text = line_text
                    run.font.size = Pt(max(6, min(line_size, 36)))
                    run.font.bold = line_bold
                    run.font.italic = line_italic
                    if line_color:
                        run.font.color.rgb = line_color

        pdf_document.close()
        prs.save(output_path)

        duration = round(time.time() - start_time, 2)
        print(f"Conversion completed in {duration} seconds")

        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"Output file created: {output_path} ({file_size} bytes)")
            return True
        else:
            print("Error: Output file not created")
            return False

    except Exception as e:
        print(f"Conversion error: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pdf_to_powerpoint_converter.py <input_pdf> <output_pptx>")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]

    if not os.path.exists(input_file):
        print(f"Error: Input file not found: {input_file}")
        sys.exit(1)

    success = convert_pdf_to_powerpoint(input_file, output_file)
    sys.exit(0 if success else 1)
